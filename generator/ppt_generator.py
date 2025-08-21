from pathlib import Path
from pptx import Presentation

# Folder containing images
FOLDER = Path("samples/")
IMAGE_EXTENSIONS = [".png", ".jpg", ".jpeg"]


def load_background():
    """Load the backgrounds to be used in the presentation"""
    images = [str(f) for f in FOLDER.iterdir() if f.suffix.lower() in IMAGE_EXTENSIONS]
    if len(images) == 0:
        print("No images were found")
        return None
    return images

prs = Presentation()


def create_ppt(content, output_file):
    print("Creating ppt...")

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    images = load_background()

    print(type(slide.background))

    left = top = 0
    slide.shapes.add_picture(images[0], left-0.1*prs.slide_width, top, height = prs.slide_height)

    prs.save(output_file)
